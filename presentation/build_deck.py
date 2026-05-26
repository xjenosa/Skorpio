"""
Skorpio hackathon pitch deck generator.

Generates a 10-slide PowerPoint deck styled to match the Skorpio website
(dark theme, orange accent, mono + serif typography). Designed for a
~5-minute pitch.

Usage (from project root):
    pip install python-pptx
    python presentation/build_deck.py

Output: presentation/Skorpio-Pitch.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Brand palette (mirrors frontend/src/styles/colors_and_type.css) ─────── #
BG_0 = RGBColor(0x0A, 0x0A, 0x09)
BG_1 = RGBColor(0x16, 0x16, 0x14)
BG_2 = RGBColor(0x1F, 0x1D, 0x1A)
FG_1 = RGBColor(0xFA, 0xF9, 0xF5)
FG_2 = RGBColor(0xB3, 0xB1, 0xA8)
FG_3 = RGBColor(0x6F, 0x6C, 0x63)
ACCENT = RGBColor(0xF3, 0x87, 0x64)
ACCENT_SOFT = RGBColor(0x2A, 0x1B, 0x14)
RULE = RGBColor(0x2A, 0x28, 0x24)

# Pipeline accents (mirrors frontend/src/pipelines.ts)
PIPE_SITING = RGBColor(0xF3, 0x87, 0x64)
PIPE_EXPANSION = RGBColor(0x7A, 0xA2, 0xF7)
PIPE_WINTER = RGBColor(0x6F, 0xCF, 0x8E)
PIPE_ELECTRIFY = RGBColor(0xE0, 0xA4, 0x4A)
PIPE_INVESTMENT = RGBColor(0xC8, 0x89, 0xE6)

# Universally available fonts so the deck renders consistently on any
# presenter machine (no need to install Inter / JetBrains Mono).
FONT_SANS = "Calibri"
FONT_MONO = "Consolas"
FONT_SERIF = "Georgia"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────── #

def _zero_margins(text_frame):
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0


def add_bg(slide, color=BG_0):
    """Full-bleed background rectangle (slide.background API is unreliable)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_rect(slide, left, top, width, height, *, fill=BG_1, border=None, border_color=RULE):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if border is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = border_color
        rect.line.width = Pt(border)
    rect.shadow.inherit = False
    return rect


def add_dot(slide, left, top, color, size_pt=10):
    """Colored circle. size_pt is the diameter in points."""
    side = Pt(size_pt)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, side, side)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return dot


def add_text(slide, text, left, top, width, height, *,
             size=14, font=FONT_SANS, color=FG_1, bold=False,
             italic=False, align=PP_ALIGN.LEFT, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_runs(slide, segments, left, top, width, height, *, align=PP_ALIGN.LEFT,
             line_spacing=1.2):
    """A single paragraph with multiple styled runs.
    segments: list of dicts {text, size, color, bold, italic, font}.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg in segments:
        run = p.add_run()
        run.text = seg["text"]
        run.font.name = seg.get("font", FONT_SANS)
        run.font.size = Pt(seg.get("size", 14))
        run.font.bold = seg.get("bold", False)
        run.font.italic = seg.get("italic", False)
        run.font.color.rgb = seg.get("color", FG_1)
    return box


def add_bullets(slide, items, left, top, width, height, *,
                size=16, lead_color=FG_1, body_color=FG_2,
                glyph="—", glyph_color=ACCENT, space_after=10):
    """Bullet list. items: each item is a string OR (lead, body) tuple."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.25

        glyph_run = p.add_run()
        glyph_run.text = f"{glyph}  "
        glyph_run.font.name = FONT_SANS
        glyph_run.font.size = Pt(size)
        glyph_run.font.color.rgb = glyph_color

        if isinstance(item, tuple):
            lead, body = item
            lead_run = p.add_run()
            lead_run.text = lead
            lead_run.font.name = FONT_SANS
            lead_run.font.size = Pt(size)
            lead_run.font.bold = True
            lead_run.font.color.rgb = lead_color

            body_run = p.add_run()
            body_run.text = " " + body
            body_run.font.name = FONT_SANS
            body_run.font.size = Pt(size)
            body_run.font.color.rgb = body_color
        else:
            run = p.add_run()
            run.text = item
            run.font.name = FONT_SANS
            run.font.size = Pt(size)
            run.font.color.rgb = lead_color
    return box


# ── Composite helpers (consistent slide chrome) ────────────────────────── #

def add_wordmark(slide, *, left=Inches(0.6), top=Inches(0.45), size=12):
    """SK[O]RPIO wordmark — the O in accent."""
    box = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.4))
    tf = box.text_frame
    _zero_margins(tf)
    p = tf.paragraphs[0]
    for letters, color in [("SK", FG_1), ("O", ACCENT), ("RPIO", FG_1)]:
        run = p.add_run()
        run.text = letters
        run.font.name = FONT_SANS
        run.font.bold = True
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_eyebrow(slide, text, *, left=Inches(0.6), top=Inches(1.2), size=11):
    return add_text(slide, text.upper(), left, top, Inches(10), Inches(0.4),
                    size=size, font=FONT_MONO, color=FG_3)


def add_title(slide, text, *, left=Inches(0.6), top=Inches(1.7), width=Inches(11.5),
              size=44, color=FG_1):
    return add_text(slide, text, left, top, width, Inches(1.4),
                    size=size, font=FONT_SERIF, color=color, line_spacing=1.05)


def add_page_number(slide, n, total, *, top=Inches(7.0)):
    add_text(slide, f"{n:02d} / {total:02d}", Inches(12.4), top, Inches(0.7),
             Inches(0.3), size=9, font=FONT_MONO, color=FG_3)


def add_chrome(slide, n, total, eyebrow):
    """Apply the consistent header chrome to a content slide."""
    add_bg(slide)
    add_wordmark(slide)
    add_eyebrow(slide, eyebrow)
    add_page_number(slide, n, total)


# ── Slide factories ────────────────────────────────────────────────────── #

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    # Subtle background grid (single soft accent dot, top-right hint)
    add_dot(s, Inches(12.4), Inches(0.55), ACCENT, size_pt=8)
    add_text(s, "GRID · LIVE", Inches(11.1), Inches(0.5), Inches(1.2),
             Inches(0.3), size=10, font=FONT_MONO, color=FG_3,
             align=PP_ALIGN.RIGHT)

    add_wordmark(s, size=14)
    add_eyebrow(s, "PLANNER DASHBOARD · V0.2", top=Inches(2.1))

    add_text(s, "Reason about the grid",
             Inches(0.6), Inches(2.7), Inches(12), Inches(1.4),
             size=68, font=FONT_SERIF, color=FG_1, line_spacing=1.05)
    add_text(s, "from one prompt.",
             Inches(0.6), Inches(3.85), Inches(12), Inches(1.4),
             size=68, font=FONT_SERIF, color=ACCENT,
             italic=True, line_spacing=1.05)

    add_text(s, "An AI co-pilot for the people who plan Canada's electricity grid.",
             Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             size=20, color=FG_2)

    add_text(s, "SENECA HACKATHON · SMART GRID RESILIENCE & PREPARING FOR ELECTRIFICATION",
             Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
             size=10, font=FONT_MONO, color=FG_3)


def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "01 · The problem")
    add_title(s, "Canada's grid is being asked to do more than ever.")

    add_bullets(s, [
        ("Electric cars and heat pumps",
         "are changing how much power we use, faster than utilities can plan for."),
        ("Extreme weather",
         "like deep-freezes and heatwaves puts the grid — and people — at risk."),
        ("Big grid decisions",
         "still take months, cost millions, and rely on outside consultants."),
        ("Today's tools are stuck",
         "in spreadsheets — planners juggle weather, maps, and grid data by hand."),
    ], Inches(0.6), Inches(3.4), Inches(12), Inches(3.5), size=18)


def slide_solution(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "02 · The solution")
    add_title(s, "One question. Five tools. The AI picks the right one.")

    add_bullets(s, [
        ("Just type what you need —", "no settings, no forms, no jargon."),
        ("An AI assistant", "picks the right kind of analysis for your question."),
        ("Watch it think live —", "every step, every data source, in real time."),
        ("Powered by Claude AI", "with custom tools we built for energy data."),
        ("Easy to grow —", "new analysis types can be added without rebuilding the app."),
    ], Inches(0.6), Inches(3.4), Inches(12), Inches(3.5), size=17)


def slide_capabilities(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "03 · Capabilities")
    add_title(s, "Five tools, one assistant.")

    pipelines = [
        ("Datacenter Siting", PIPE_SITING,
         "Finds the best place in Canada to build a new AI data center, balancing carbon intensity, latency, and cost."),
        ("Expansion Planner", PIPE_EXPANSION,
         "Figures out how to grow an existing data center site."),
        ("Winter Peak Stress Tester", PIPE_WINTER,
         "Simulates a deep cold snap to see how heat pumps, EVs, and electric heaters push the grid."),
        ("Electrification Readiness", PIPE_ELECTRIFY,
         "Tells you which neighborhoods are ready for more electric vehicles and heat pumps, and which need help first."),
        ("Grid Investment Optimizer", PIPE_INVESTMENT,
         "Takes a fixed upgrade budget and finds the smartest way to spend it for the most protection against future climate disasters."),
    ]

    row_top = Inches(3.2)
    row_h = Inches(0.72)
    row_gap = Inches(0.08)
    row_w = Inches(12.1)
    for i, (name, color, blurb) in enumerate(pipelines):
        top = row_top + (row_h + row_gap) * i
        add_rect(s, Inches(0.6), top, row_w, row_h, fill=BG_1)
        # Color stripe on the left edge
        add_rect(s, Inches(0.6), top, Inches(0.06), row_h, fill=color)
        add_dot(s, Inches(0.85), top + Inches(0.27), color, size_pt=12)
        add_text(s, name, Inches(1.18), top + Inches(0.15), Inches(4.4),
                 Inches(0.45), size=15, color=color, bold=True)
        add_text(s, blurb, Inches(1.18), top + Inches(0.4), Inches(10.8),
                 Inches(0.4), size=12, color=FG_2)


def slide_hackathon_fit(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "04 · Challenge fit")
    add_title(s, "Built for the challenge.")

    add_text(s,
             "Smart Grid Resilience & Preparing for Electrification",
             Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
             size=13, font=FONT_MONO, color=FG_3, italic=True)

    # Two columns: Resilience | Electrification
    col_w = Inches(5.95)
    col_h = Inches(3.4)

    # Resilience card
    add_rect(s, Inches(0.6), Inches(3.6), col_w, col_h, fill=BG_1)
    add_text(s, "RESILIENCE", Inches(0.85), Inches(3.8), col_w, Inches(0.4),
             size=11, font=FONT_MONO, color=ACCENT)
    add_text(s, "When the grid is stressed.",
             Inches(0.85), Inches(4.15), col_w, Inches(0.5),
             size=20, font=FONT_SERIF, color=FG_1)
    add_bullets(s, [
        ("Winter Peak Stress Tester", "— see how a deep-freeze hits your grid"),
        ("Investment Optimizer", "— spend upgrade money where it protects most"),
    ], Inches(0.85), Inches(5.0), Inches(5.5), Inches(2.0), size=13)

    # Electrification card
    add_rect(s, Inches(6.78), Inches(3.6), col_w, col_h, fill=BG_1)
    add_text(s, "ELECTRIFICATION", Inches(7.03), Inches(3.8), col_w, Inches(0.4),
             size=11, font=FONT_MONO, color=ACCENT)
    add_text(s, "When demand is changing.",
             Inches(7.03), Inches(4.15), col_w, Inches(0.5),
             size=20, font=FONT_SERIF, color=FG_1)
    add_bullets(s, [
        ("Electrification Readiness", "— score neighborhoods for EV and heat pump growth"),
        ("Datacenter Siting + Expansion", "— place big new power users wisely"),
    ], Inches(7.03), Inches(5.0), Inches(5.5), Inches(2.0), size=13)

    add_text(s, "Canadian-first — uses live data from Ontario (IESO), Alberta (AESO), Quebec, and Natural Resources Canada",
             Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             size=11, font=FONT_MONO, color=FG_3)


def slide_market(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "05 · Market opportunity")
    add_title(s, "$1 trillion+ in Canadian grid investment by 2050.")

    # Big number tile
    add_rect(s, Inches(0.6), Inches(3.4), Inches(4.0), Inches(3.6), fill=BG_1)
    add_text(s, "BY 2050", Inches(0.85), Inches(3.6), Inches(3.5),
             Inches(0.4), size=10, font=FONT_MONO, color=ACCENT)
    add_text(s, "$1T+", Inches(0.85), Inches(4.0), Inches(3.5), Inches(1.4),
             size=66, font=FONT_SERIF, color=FG_1)
    add_text(s,
             "Cost of doubling Canada's electricity grid to power EVs, heat pumps, and net-zero goals.",
             Inches(0.85), Inches(5.6), Inches(3.5), Inches(1.2),
             size=12, color=FG_2)
    add_text(s,
             "Source: NRCan, Powering Canada Strong (national electricity strategy, May 2026)",
             Inches(0.85), Inches(6.7), Inches(3.5), Inches(0.3),
             size=9, font=FONT_MONO, color=FG_3)

    # Right column: who could buy this, in plain English
    add_bullets(s, [
        ("Whole market:", "software for planning Canada's electricity grid."),
        ("Realistic share:", "60+ Canadian utilities, plus federal and provincial planners and big cities."),
        ("First customers:", "mid-size local utilities facing the EV/heat pump wave first."),
        ("Adjacent market:", "global grid planning software, $2.04B in 2025 growing to $4.12B by 2030 (15.2% CAGR)."),
    ], Inches(5.0), Inches(3.5), Inches(7.6), Inches(3.2), size=15)

    add_text(s,
             "Sources: NRCan Powering Canada Strong, 2026 (grid investment); The Business Research Company, Grid Planning Software Global Market Report, 2026 (adjacent market).",
             Inches(5.0), Inches(6.85), Inches(7.6), Inches(0.4),
             size=8, font=FONT_MONO, color=FG_3)


def slide_revenue(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "06 · Revenue model")
    add_title(s, "Three ways customers pay.")
    add_text(s,
             "Pricing benchmarks based on comparable utility-planning tools (PowerWorld, Energy Exemplar AURORA, Siemens PSS/E). Final pricing TBD after pilots.",
             Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
             size=10, font=FONT_MONO, color=FG_3, italic=True)

    streams = [
        ("Monthly subscriptions", "Utility planners — about $1,200 per seat per month",
         "Recurring", PIPE_SITING),
        ("Enterprise contracts", "Large utilities and federal agencies — $250K to $2M per year",
         "Recurring", PIPE_EXPANSION),
        ("One-off studies", "Consultants and developers — $25K to $100K per report",
         "Transactional", PIPE_WINTER),
        ("Premium data add-ons", "Specialized datasets and live grid feeds",
         "Add-on", PIPE_ELECTRIFY),
        ("Setup & training", "Hands-on help getting big customers up and running",
         "Services", PIPE_INVESTMENT),
    ]

    row_top = Inches(3.4)
    row_h = Inches(0.7)
    gap = Inches(0.08)
    for i, (name, desc, kind, color) in enumerate(streams):
        top = row_top + (row_h + gap) * i
        add_rect(s, Inches(0.6), top, Inches(12.1), row_h, fill=BG_1)
        add_dot(s, Inches(0.85), top + Inches(0.27), color, size_pt=10)
        add_text(s, name, Inches(1.15), top + Inches(0.18), Inches(3.6),
                 Inches(0.4), size=14, color=FG_1, bold=True)
        add_text(s, desc, Inches(4.85), top + Inches(0.2), Inches(6.0),
                 Inches(0.4), size=12, color=FG_2)
        add_text(s, kind.upper(), Inches(11.0), top + Inches(0.22), Inches(1.6),
                 Inches(0.4), size=10, font=FONT_MONO, color=color,
                 align=PP_ALIGN.RIGHT)


def slide_roadmap(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "07 · Roadmap")
    add_title(s, "From hackathon to platform.")

    phases = [
        ("NOW",  "Q2 '26", "First version",
         "5 analysis types\nFirst customer testing\nWeb dashboard", ACCENT),
        ("NEXT", "Q3 '26", "Live grid data",
         "Real-time feeds\nMultiple companies at once\nFull activity logs", PIPE_EXPANSION),
        ("LATER","Q4 '26", "Plug-in mode",
         "Connects to grid mapping\n  and control software\nCustomer-branded versions", PIPE_WINTER),
        ("2027", "2027",   "Self-driving",
         "AI-improved forecasting\nAuto-generated what-if scenarios", PIPE_INVESTMENT),
    ]

    col_w = Inches(2.95)
    gap = Inches(0.15)
    top = Inches(3.4)
    h = Inches(3.5)

    for i, (label, when, title, body, color) in enumerate(phases):
        left = Inches(0.6) + (col_w + gap) * i
        add_rect(s, left, top, col_w, h, fill=BG_1)
        # Top stripe
        add_rect(s, left, top, col_w, Inches(0.06), fill=color)
        add_text(s, label, left + Inches(0.25), top + Inches(0.25), col_w,
                 Inches(0.35), size=10, font=FONT_MONO, color=color)
        add_text(s, when, left + Inches(0.25), top + Inches(0.55), col_w,
                 Inches(0.35), size=10, font=FONT_MONO, color=FG_3)
        add_text(s, title, left + Inches(0.25), top + Inches(0.95), col_w,
                 Inches(0.6), size=22, font=FONT_SERIF, color=FG_1)
        # Multi-line body
        box = s.shapes.add_textbox(left + Inches(0.25), top + Inches(1.7),
                                   col_w - Inches(0.5), Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True
        _zero_margins(tf)
        for j, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = line
            run.font.name = FONT_SANS
            run.font.size = Pt(13)
            run.font.color.rgb = FG_2


def slide_whats_next(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "08 · What's next")
    add_title(s, "Next 90 days.")

    add_bullets(s, [
        ("Pilot with 1–2 Canadian utilities", "to prove the workflow on real grid data."),
        ("Build a Remote Community tool", "to help Canada's ~280 off-grid towns switch from diesel to renewables."),
        ("Live grid stress dashboard", "showing real-time conditions and what people can do to help."),
        ("Open the platform", "so outside teams can build their own analysis tools on Skorpio."),
        ("Grow the team", "— add 1 backend engineer, 1 data engineer, 1 customer success lead."),
    ], Inches(0.6), Inches(3.4), Inches(12), Inches(3.7), size=16)


def slide_tech_stack(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(s, n, total, "09 · Tech stack")
    add_title(s, "Built modern, ready to grow, Canadian data first.")

    blocks = [
        ("WHAT YOU SEE", FG_2,
         "A clean web dashboard\n(React, TypeScript)\nbuilt for grid planners"),
        ("BEHIND THE SCENES", FG_2,
         "A Python web service\nstreams answers live\nas the AI reasons"),
        ("THE BRAIN", ACCENT,
         "Claude AI\n(Sonnet, Opus, Haiku)\nwith custom energy tools"),
        ("THE DATA", PIPE_WINTER,
         "Live grid feeds:\nIESO, AESO, Hydro-Québec\nWeather, prices, carbon"),
        ("HOW IT RUNS", PIPE_EXPANSION,
         "Docker containers,\nSnowflake for analytics,\nautomated testing"),
    ]

    col_w = Inches(2.36)
    gap = Inches(0.14)
    top = Inches(3.4)
    h = Inches(3.4)

    for i, (label, color, body) in enumerate(blocks):
        left = Inches(0.6) + (col_w + gap) * i
        add_rect(s, left, top, col_w, h, fill=BG_1)
        add_text(s, label, left + Inches(0.25), top + Inches(0.25), col_w,
                 Inches(0.35), size=10, font=FONT_MONO, color=color)
        # Multi-line body
        box = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.85),
                                   col_w - Inches(0.5), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        _zero_margins(tf)
        for j, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = line
            run.font.name = FONT_SANS
            run.font.size = Pt(13)
            run.font.color.rgb = FG_1

    add_text(s,
             "Thanks — questions?",
             Inches(0.6), Inches(7.0), Inches(11), Inches(0.4),
             size=14, font=FONT_SERIF, color=ACCENT, italic=True)


# ── Build ──────────────────────────────────────────────────────────────── #

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    factories = [
        slide_problem,
        slide_solution,
        slide_capabilities,
        slide_hackathon_fit,
        slide_market,
        slide_revenue,
        slide_roadmap,
        slide_whats_next,
        slide_tech_stack,
    ]
    total = len(factories)  # title slide is unnumbered

    slide_title(prs)
    for i, factory in enumerate(factories, start=1):
        factory(prs, i, total)

    out_dir = Path(__file__).resolve().parent
    out_path = out_dir / "Skorpio-Pitch.pptx"
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
